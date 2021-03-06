import logging
import tempfile

import numpy as np

from pptx import Presentation
from pptx.util import Pt


def _iter_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell


def _animate(axes, width=300, height=200):
    images = []
    for fig, ax in axes:
        fig.canvas.draw()  # draw the canvas, cache the renderer
        image = np.frombuffer(fig.canvas.tostring_rgb(), dtype='uint8')
        image = image.reshape(fig.canvas.get_width_height()[::-1] + (3,))
        # image.reshape((width, height))
        images.append(image)

    return images


class pptx_image:
    """
        A wrapper for an image filename. It is to be used when wanting so insert an image from a file
    """
    def __init__(self, filename):
        self.filename = filename


class Slideshow:
    """
        Class for creating PPTX slideshows. It must be initialized with the path to template.
        Aferwards, the Execute method compiles a pptx presentaiton, and the SaveTo method writes it to a file.
    """
    def __init__(self, template_path,
                 string2colormap={},
                 dpi=200):
        """
        Object that wraps the creation of slideshows and exposes functions that are useful for datascience outputs.

        :param template_path: The path to the PPTX template
        :param string2colormap: a dictionary whose keys are strings, and whose values are pptx.dml.color.RGBColor objects.
        When writing a table, cells whose value matches a key in this dictionary, will have font size equal to its value.
        :param dpi: Resolution to write images in
        """
        self.dpi = dpi
        self.slideshow = Presentation(template_path)
        self.string2colormap = string2colormap

    def SaveTo(self, path):
        """
        Saves to file
        :param path: Path were to save the resulting PPTX slideshow.
        """
        self.slideshow.save(path)

    def Execute(self, slides):
        """
        Compiles a set of slides.

        :param slides: a list of dictionaries, each of which describes the content of a slide. One example would be
        ```          {
            "name": "Slide3",
            "layout": "LayoutName",
            'placeholders': {
                "Image Place Holder": matplotlib_figure_object
            } ```
            Where ```layout``` is the name of a Slide Layout, and each key in Placeholders is an object placeholder.
            These are defined in the PPTX template with which the object was initialized.
            The values of the ```placeholders``` disctionaries are objects of a set of supported classes.
            Currently, the following are supported:
                matplotlib.Figure : will be rendered as an image
                data2pptx.pptx_image : will be rendered as an image
                string : to be rendered as text
                pandas.dataframe : to be rendered as a table
                bytes : will be rendered as an image. Useful for animated gifs.
        """
        slds = self.slideshow
        for sl in slides:
            layout = self.FindLayout(sl['layout'])
            if layout:
                slide = slds.slides.add_slide(layout)
            else:
                raise Exception('Slide layout: {} does not exist in this template.' +
                                ' the available layouts are: {}'.format(sl['layout'],
                                                                        ' -, - '.join(
                                                                            [i.name for i in
                                                                             slds.slide_layouts])))

            if 'title' in sl:
                title_shape = slide.shapes.title
                if title_shape:
                    title_shape.text = sl['title']

            if 'placeholders' in sl:
                for k in sl['placeholders']:
                    pl = self.FindPlaceholder(slide, k)
                    if pl:
                        ty = type(sl['placeholders'][k]).__name__.lower()
                        method_name = "_put_" + ty
                        try:
                            method = getattr(self, method_name)
                        except:
                            raise Exception(
                                'There is no function to add an object of type {} onto placeholder {}'.format(ty, k))
                        method(pl, sl['placeholders'][k])
                    else:
                        raise Exception('Placeholder {}, is unknown in this slide-type. ' +
                                        'Available placeholders are: {}'.format(k,
                                                                                ', '.join(
                                                                                    [
                                                                                        i.name for i in
                                                                                        slide.placeholders])))

    def FindPlaceholder(self, slide, name):
        """

        :param slide:
        :param name:
        :return:
        """
        return next((i for i in slide.placeholders if i.name == name), None)

    def FindLayout(self, name):
        """

        :param name:
        :return:
        """
        return next((i for i in self.slideshow.slide_layouts if i.name == name), None)

    def GetLayouts(self):
        """

        :return:
        """
        return [sl for sl in self.slideshow.slide_layouts]

    def _put_str(self, placeholder, text):
        placeholder.text = text

    def _put_figure(self, placeholder, fig):
        logging.debug(" ".join([str(x) for x in ["figure size (figure):",
                                                 placeholder.width,
                                                 placeholder.height,
                                                 (placeholder.width / placeholder.height)]]))
        fp = tempfile.TemporaryFile()
        fig.savefig(fp, dpi=self.dpi, format="png")
        fp.seek(0)
        placeholder.insert_picture(fp)

    def _put_pptx_image(self, placeholder, image):
        placeholder.insert_picture(image.filename)

    def _put_dataframe(self, placeholder, dataframe):
        x = placeholder.insert_table(len(dataframe) + 1, len(dataframe.columns) + len(dataframe.index.names))
        table = x.table
        table.columns[1].width = Pt(50)
        table.columns[0].width = Pt(150)
        table.columns[2].width = Pt(50)
        table.columns[3].width = Pt(50)

        rowId = 0
        colId = 0

        # if Indices don't have names, we replace them by blanks
        empty_indices = [i for i, x in enumerate(dataframe.index.names) if x is None]
        if len(empty_indices) > 0:
            allindices = list(dataframe.index.names)
            for ii in empty_indices:
                allindices[ii] = " "
            dataframe.index.names = allindices
        for row in dataframe.itertuples():
            if rowId == 0:
                for col in dataframe.index.names:
                    table.cell(0, colId).text = col
                    colId = colId + 1

                for col in dataframe.columns:
                    table.cell(0, colId).text = col
                    colId = colId + 1

                rowId = 1

            colId = 0
            for col in row:
                table.cell(rowId, colId).text = str(col)
                if col in self.string2colormap.keys():
                    paragraph = table.cell(rowId, colId).text_frame.paragraphs[0]
                    for run in paragraph.runs:
                        run.font.color.rgb = self.string2colormap[col]
                colId = colId + 1

            rowId = rowId + 1

        for cell in _iter_cells(table):
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(12)

    def _put_bytes(self, placeholder, the_bytes):
        logging.debug(" ".join([str(x) for x in ["figure size (bytes):",
                                                 placeholder.width,
                                                 placeholder.height,
                                                 (placeholder.width / placeholder.height)]]))
        # path = "/tmp/" + str(uuid.uuid4()) + ".gif"
        # with open(path, "wb") as fout:
        #     fout.write(the_bytes)
        # placeholder.insert_picture(path)
        fp = tempfile.TemporaryFile()
        fp.write(the_bytes)
        fp.seek(0)
        placeholder.insert_picture(fp)
